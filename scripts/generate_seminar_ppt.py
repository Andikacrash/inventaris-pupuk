# -*- coding: utf-8 -*-
"""Generate PPT Seminar Pendadaran — Andika Suyandra (struktur mengikuti contoh)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(r"d:\SKRIPSI\Skripsi\PPT_Seminar_Pendadaran_Andika_Suyandra.pptx")

# Warna tema (akademik biru, mirip contoh)
C_NAVY = RGBColor(0x1B, 0x3A, 0x6B)
C_BLUE = RGBColor(0x2E, 0x6B, 0xB5)
C_DARK = RGBColor(0x1E, 0x17, 0x10)
C_MUTED = RGBColor(0x4A, 0x37, 0x28)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT = RGBColor(0x1A, 0x5C, 0x42)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title, section=""):
    """Header biru di atas slide isi."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    if section:
        p2 = tf.add_paragraph()
        p2.text = section
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)


def add_bullets(slide, items, left=0.6, top=1.3, width=12.0, height=5.5, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = C_DARK
        p.space_after = Pt(8)


def add_two_col(slide, left_title, left_items, right_title, right_items):
    add_title_bar(slide, left_title.split(" — ")[0] if " — " in left_title else left_title)
    # left col title
    lbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.5))
    ltf = lbox.text_frame
    ltf.text = left_title
    ltf.paragraphs[0].font.size = Pt(22)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = C_BLUE
    add_bullets(slide, left_items, left=0.5, top=1.8, width=5.8, size=18)

    rbox = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(5.8), Inches(0.5))
    rtf = rbox.text_frame
    rtf.text = right_title
    rtf.paragraphs[0].font.size = Pt(22)
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.color.rgb = C_BLUE
    add_bullets(slide, right_items, left=6.8, top=1.8, width=5.8, size=18)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_NAVY)

    # Judul
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "RANCANG BANGUN SISTEM INFORMASI PENJUALAN PUPUK DENGAN METODE HUMAN-CENTERED DESIGN"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "(STUDI KASUS: TOKO PUPUK SAWIJI TANI)"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)

    # Nama
    nbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.8))
    nf = nbox.text_frame
    np = nf.paragraphs[0]
    np.text = "ANDIKA SUYANDRA"
    np.font.size = Pt(24)
    np.font.bold = True
    np.font.color.rgb = C_WHITE
    np.alignment = PP_ALIGN.CENTER

    np2 = nf.add_paragraph()
    np2.text = "2200018237"
    np2.font.size = Pt(18)
    np2.font.color.rgb = C_WHITE
    np2.alignment = PP_ALIGN.CENTER

    # Pembimbing & penguji (placeholder — sesuaikan jika sudah fix)
    pbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8))
    pf = pbox.text_frame
    for line in [
        "Jefree Fahana, S.T., M.Kom.  —  Dosen Pembimbing",
        "Ali Tarmuji, S.T., M.Cs.  —  Dosen Penguji 1",
        "Bambang Robi'in, S.T., M.T.  —  Dosen Penguji 2",
    ]:
        pp = pf.paragraphs[0] if line == pf.paragraphs[0].text or not pf.paragraphs[0].text else pf.add_paragraph()
        if pp.text == "":
            pp.text = line
        else:
            pnew = pf.add_paragraph()
            pnew.text = line
        pp.font.size = Pt(14)
        pp.font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)
        pp.alignment = PP_ALIGN.CENTER


def slide_toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    title.text_frame.text = "Pembahasan"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = C_NAVY

    items = [
        ("1", "Pendahuluan"),
        ("2", "Tinjauan Pustaka"),
        ("3", "Metodologi Penelitian"),
        ("4", "Hasil dan Pembahasan"),
        ("5", "Kesimpulan dan Saran"),
    ]
    y = 1.6
    for num, label in items:
        circle = slide.shapes.add_shape(9, Inches(1.0), Inches(y), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = C_BLUE
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.text = num
        ctf.paragraphs[0].font.size = Pt(18)
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].font.color.rgb = C_WHITE
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

        lbl = slide.shapes.add_textbox(Inches(1.8), Inches(y), Inches(10), Inches(0.55))
        lbl.text_frame.text = label
        lbl.text_frame.paragraphs[0].font.size = Pt(22)
        lbl.text_frame.paragraphs[0].font.color.rgb = C_DARK
        y += 0.85


def slide_content(prs, header, section, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, header, section)
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12), Inches(0.4))
        sbox.text_frame.text = subtitle
        sbox.text_frame.paragraphs[0].font.size = Pt(20)
        sbox.text_frame.paragraphs[0].font.bold = True
        sbox.text_frame.paragraphs[0].font.color.rgb = C_BLUE
        top = 1.65
    else:
        top = 1.3
    add_bullets(slide, bullets, top=top)


def slide_placeholder_image(prs, header, caption):
    """Slide untuk gambar (Use Case, ERD, UI) — sisipkan screenshot dari skripsi."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, header)

    rect = slide.shapes.add_shape(1, Inches(1.5), Inches(1.4), Inches(10.3), Inches(4.8))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xF4, 0xF1, 0xEC)
    rect.line.color.rgb = C_BLUE
    rect.line.width = Pt(2)

    tf = rect.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"[Sisipkan gambar: {caption}]"
    p.font.size = Pt(20)
    p.font.color.rgb = C_MUTED
    p.alignment = PP_ALIGN.CENTER

    note = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12), Inches(0.5))
    note.text_frame.text = "Ambil dari Bab IV skripsi atau screenshot aplikasi lokal."
    note.text_frame.paragraphs[0].font.size = Pt(12)
    note.text_frame.paragraphs[0].font.color.rgb = C_MUTED


def slide_sus_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "Pengujian", "System Usability Scale (SUS)")

    rows, cols = 7, 3
    tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.5)).table
    headers = ["No", "Responden", "Skor SUS"]
    data = [
        ["1", "Responden 1", "82,5"],
        ["2", "Responden 2", "87,5"],
        ["3", "Responden 3", "80,0"],
        ["4", "Responden 4", "82,5"],
        ["5", "Responden 5", "87,5"],
        ["Rata-rata", "", "84,0"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = C_WHITE
            p.font.bold = True
            p.font.size = Pt(14)
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = C_DARK

    add_bullets(
        slide,
        [
            "Rata-rata skor SUS: 84,0 — kategori Acceptable / Excellent.",
            "Antarmuka Warm Neutral (HCD) mudah dipahami pemilik, admin, dan kasir.",
            "Pengguna dapat menyelesaikan tugas uji tanpa bantuan berulang.",
        ],
        top=4.2,
        size=18,
    )


def slide_blackbox_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "Pengujian", "Black Box Testing")

    rows, cols = 8, 4
    tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.2)).table
    headers = ["No", "Modul", "Jumlah Skenario", "Hasil"]
    modul = [
        ["1", "Login & Akun Pengguna", "10", "Valid"],
        ["2", "Hak Akses (Admin/Kasir)", "10", "Valid"],
        ["3", "Kelola Barang & Supplier", "16", "Valid"],
        ["4", "Kasir (POS)", "10", "Valid"],
        ["5", "Hutang & Cicilan", "17", "Valid"],
        ["6", "Laporan, Dashboard, Stok", "35", "Valid"],
        ["", "TOTAL", "98", "100% Valid"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = C_WHITE
            p.font.bold = True
            p.font.size = Pt(13)
    for r, row in enumerate(modul, start=1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val
            for p in tbl.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(12)

    add_bullets(
        slide,
        [
            "98 skenario black box — seluruh fungsi sesuai spesifikasi (100%).",
            "Validasi input, konsistensi stok, hak akses, dan laporan berjalan baik.",
        ],
        top=4.7,
        size=17,
    )


def slide_tum(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "Pengujian", "Task Usability Metrics (Efisiensi Kerja)")

    rows, cols = 4, 4
    tbl = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.3), Inches(12.5), Inches(2.8)).table
    headers = ["Skenario", "Sebelum (Manual)", "Sesudah (Sistem)", "Dampak"]
    data = [
        [
            "Transaksi & kuitansi",
            "2–4 menit; error 30%",
            "30–45 detik; error 0%",
            "Waktu ↓ ~80%",
        ],
        [
            "Cek stok barang",
            "2–5 menit; error 25%",
            "3–5 detik; error 0%",
            "Stok real-time",
        ],
        [
            "Laporan bulanan",
            "3–4 hari; error 40%",
            "< 30 detik; error 0%",
            "Rekap instan",
        ],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = C_WHITE
            p.font.bold = True
            p.font.size = Pt(12)
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_toc(prs)

    # --- BAB I Pendahuluan ---
    slide_content(
        prs,
        "Pendahuluan",
        "",
        [
            "Latar Belakang",
            "• Pencatatan penjualan & stok di Toko Pupuk Sawiji Tani masih manual/semi-digital.",
            "• Volume ~Rp500 juta/bulan, >1.000 transaksi/bulan, 30–50 transaksi/hari.",
            "",
            "Permasalahan",
            "• Kesalahan pencatatan, duplikasi data, nota hilang, laporan lambat (hingga 4 hari).",
            "• Stok tidak real-time → miskomunikasi dengan pelanggan.",
            "",
            "Solusi",
            "• Sistem Informasi Inventaris & Penjualan Pupuk berbasis web (Laravel + HCD).",
            "",
            "Tujuan",
            "• Digitalisasi transaksi, stok, piutang, laporan dengan antarmuka ramah pengguna dewasa.",
        ],
    )

    slide_content(
        prs,
        "Pendahuluan",
        "Batasan Masalah",
        [
            "Fokus pada pengembangan SIM transaksi penjualan & pengelolaan stok Toko Pupuk Sawiji Tani.",
            "Menggunakan Laravel dan pendekatan Human-Centered Design (ISO 9241-210).",
            "Sistem berbasis web, akses internal admin & kasir.",
            "Tanpa integrasi pembayaran digital atau e-commerce.",
        ],
    )

    slide_content(
        prs,
        "Pendahuluan",
        "Rumusan Masalah",
        [
            "1. Bagaimana merancang dan mengimplementasikan sistem informasi manajemen transaksi menggunakan Laravel dan metode Human-Centered Design (HCD)?",
            "2. Bagaimana melakukan validasi sistem melalui black box testing dan user satisfaction survey (SUS & Task Usability Metrics)?",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_two_col(
        slide,
        "Tujuan Penelitian",
        [
            "1. Menghasilkan rancangan & implementasi SIM transaksi, stok, dan laporan (Laravel + HCD).",
            "2. Menghasilkan data validasi melalui black box testing dan survei kepuasan pengguna.",
        ],
        "Manfaat Penelitian",
        [
            "Praktis: digitalisasi transaksi real-time, kurangi kesalahan, laporan < 30 menit.",
            "Praktis: efisiensi operasional & akurasi stok Toko Pupuk Sawiji Tani.",
            "Akademis: referensi SI berbasis Laravel + HCD untuk UKM pertanian.",
        ],
    )

    # --- BAB II ---
    slide_content(
        prs,
        "Tinjauan Pustaka",
        "Kajian Penelitian Terdahulu",
        [
            "Dhamara et al. — SI penjualan properti + HCD; SUS 90,625 (Excellent).",
            "Azizunhakim et al. — UX web-store kopi + HCD; SUS 68 (baik).",
            "Isnanto et al. — Dashboard deteksi hoaks + HCD; Black Box & SUS 76.",
            "Penelitian ini: SI inventaris & POS pupuk + HCD, Black Box, SUS, Task Usability Metrics.",
        ],
    )

    slide_content(
        prs,
        "Tinjauan Pustaka",
        "Landasan Teori",
        [
            "1. Rancang Bangun Sistem Informasi",
            "2. Aplikasi Berbasis Web",
            "3. Human-Centered Design (ISO 9241-210)",
            "4. Framework Laravel (MVC, Eloquent ORM)",
            "5. Unified Modeling Language (UML)",
            "6. Entity Relationship Diagram (ERD)",
            "7. Pengujian Black Box",
            "8. System Usability Scale (SUS)",
            "9. Task Usability Metrics",
        ],
    )

    # --- BAB III ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_two_col(
        slide,
        "Metode Pengumpulan Data — Data Primer",
        ["1. Wawancara semi-terstruktur (pemilik, admin, kasir)", "2. Observasi proses transaksi manual"],
        "Data Sekunder",
        ["1. Dokumen operasional (stok, transaksi, laporan)", "2. Literatur & penelitian terdahulu"],
    )

    slide_content(
        prs,
        "Metodologi Penelitian",
        "Tahapan Penelitian (HCD — ISO 9241-210)",
        [
            "1. Specify the user and context of use — wawancara & observasi di toko.",
            "2. Specify user requirements — kebutuhan fungsional & antarmuka.",
            "3. Produce design solutions — prototipe 3 versi UI → tema Warm Neutral.",
            "4. Evaluate designs — Black Box (98 skenario), SUS, Task Usability Metrics.",
            "Iterasi desain hingga memenuhi ambang SUS ≥ 80 dan UAT lulus.",
        ],
    )

    # --- BAB IV ---
    slide_content(
        prs,
        "Hasil dan Pembahasan",
        "Pengumpulan Data",
        [
            "Admin: kelola produk, kategori, supplier, stok, laporan, hutang.",
            "Kasir: transaksi POS cepat, scan barcode, riwayat transaksi, bayar hutang/cicilan.",
            "Sistem: stok real-time, laporan otomatis, pencarian & filter, UI sederhana.",
            "Kebutuhan HCD: huruf besar, kontras tinggi, istilah Indonesia, format angka ribuan.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "Hasil dan Pembahasan", "Analisis Kebutuhan")
    add_two_col(
        slide,
        "Pengguna Sistem",
        [
            "1. Admin — kelola master data, stok, laporan, akun kasir",
            "2. Kasir — transaksi POS, piutang, cicilan",
        ],
        "Kebutuhan Non-Fungsional",
        [
            "Keamanan & hak akses per peran",
            "Kinerja cepat & stabil",
            "Antarmuka mudah digunakan (SUS ≥ 80)",
            "Kompatibel berbagai browser",
        ],
    )
    # overlay fungsional
    add_bullets(
        slide,
        [
            "Kebutuhan Fungsional: login, CRUD produk, POS, stok, hutang/cicilan, dashboard, laporan Excel/PDF.",
        ],
        top=5.5,
        size=16,
    )

    slide_placeholder_image(prs, "Hasil dan Pembahasan", "Use Case Diagram")
    slide_placeholder_image(prs, "Hasil dan Pembahasan", "Activity Diagram — Transaksi POS")
    slide_placeholder_image(prs, "Hasil dan Pembahasan", "Entity Relationship Diagram (ERD)")

    for ui in [
        "Halaman Login",
        "Dashboard",
        "Kasir (Point of Sale)",
        "Kelola Barang",
        "Hutang Pelanggan",
        "Laporan Penjualan",
        "Laporan Stok",
    ]:
        slide_placeholder_image(prs, "User Interface", ui)

    slide_content(
        prs,
        "Pengujian",
        "",
        [
            "Black Box Testing — pengujian fungsional berdasarkan input/output tanpa melihat kode.",
            "System Usability Scale (SUS) — mengukur kemudahan & penerimaan pengguna (skala 0–100).",
            "Task Usability Metrics — mengukur waktu tugas & tingkat kesalahan sebelum vs sesudah sistem.",
        ],
    )

    slide_sus_table(prs)
    slide_blackbox_summary(prs)
    slide_tum(prs)

    slide_content(
        prs,
        "Kesimpulan dan Saran",
        "Kesimpulan",
        [
            "SIM Inventaris & Penjualan Pupuk berhasil dibangun (Laravel 12, MySQL, tema Warm Neutral/HCD).",
            "Black box 98 skenario — 100% valid; hak akses, stok, transaksi, laporan berfungsi baik.",
            "SUS rata-rata 84,0 (Excellent); Task Metrics: error 30%→0%, kuitansi 4 menit→45 detik.",
            "Sistem layak untuk UAT dan operasional harian Toko Pupuk Sawiji Tani.",
        ],
    )

    slide_content(
        prs,
        "Kesimpulan dan Saran",
        "Saran",
        [
            "Deploy ke hosting tetap & backup database rutin.",
            "Tambah audit log dan notifikasi stok menipis.",
            "Pelatihan kasir baru & rekonsiliasi stok fisik vs sistem berkala.",
            "Penelitian lanjutan: uji beban produksi & pengembangan fitur mobile.",
        ],
    )

    # Slide penutup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_NAVY)
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.5))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Terima Kasih"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Pertanyaan & Diskusi"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
    p2.alignment = PP_ALIGN.CENTER

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Berhasil: {OUT}")
    print(f"Total slide: {len(prs.slides)}")


if __name__ == "__main__":
    main()
